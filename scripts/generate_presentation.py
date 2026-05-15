"""Generate beautiful NDKTU presentation in Uzbek (.pptx) with images."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "scripts" / "_assets"
LOGO = ROOT / "frontend" / "public" / "logo.png"
OUTPUT = ROOT / "NDKTU_Platforma_Taqdimot.pptx"

BRAND = RGBColor(0x24, 0x2C, 0xBB)
BRAND_DARK = RGBColor(0x14, 0x19, 0x6E)
BRAND_LIGHT = RGBColor(0x5A, 0x64, 0xE6)
ACCENT = RGBColor(0xF9, 0x73, 0x16)
DARK = RGBColor(0x1F, 0x29, 0x37)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_image_full(slide, path, prs):
    slide.shapes.add_picture(
        str(path), 0, 0,
        width=prs.slide_width, height=prs.slide_height,
    )


def add_rect(slide, left, top, width, height, color, *, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_round_rect(slide, left, top, width, height, color, *, corner=0.05):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font=FONT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=DARK, bullet_color=ACCENT, line_spacing=12):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_spacing)
        bullet = p.add_run()
        bullet.text = "▸  "
        bullet.font.name = FONT
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = bullet_color
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_footer(slide, prs, page_num, total):
    add_rect(
        slide, 0, prs.slide_height - Inches(0.04),
        prs.slide_width, Inches(0.04),
        BRAND,
    )
    add_text(
        slide,
        Inches(0.5), prs.slide_height - Inches(0.45),
        Inches(8), Inches(0.3),
        "Navoiy Davlat Konchilik va Texnologiyalar Universiteti",
        size=10, color=GREY,
    )
    add_text(
        slide,
        prs.slide_width - Inches(1.5), prs.slide_height - Inches(0.45),
        Inches(1), Inches(0.3),
        f"{page_num} / {total}",
        size=10, color=BRAND, bold=True, align=PP_ALIGN.RIGHT,
    )


def title_slide(prs, total):
    slide = blank_slide(prs)
    add_image_full(slide, ASSETS / "hero_title.png", prs)

    # Logo top-left on dark background — use white card
    card_left = Inches(0.6)
    card_top = Inches(0.5)
    add_round_rect(
        slide, card_left, card_top, Inches(2.4), Inches(1.35),
        WHITE, corner=0.15,
    )
    slide.shapes.add_picture(
        str(LOGO),
        card_left + Inches(0.15), card_top + Inches(0.15),
        width=Inches(2.1),
    )

    # Accent bar
    add_rect(
        slide, Inches(0.7), Inches(3.0),
        Inches(1.2), Inches(0.1), ACCENT,
    )

    add_text(
        slide, Inches(0.7), Inches(3.2),
        Inches(12), Inches(1.2),
        "Talabalar Yuz Aniqlash",
        size=54, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(0.7), Inches(4.0),
        Inches(12), Inches(1.2),
        "Platformasi",
        size=54, bold=True, color=WHITE,
    )

    # Subtitle in pill
    add_round_rect(
        slide, Inches(0.7), Inches(5.1),
        Inches(8.5), Inches(0.6),
        BRAND_DARK, corner=0.4,
    )
    add_text(
        slide, Inches(0.7), Inches(5.18),
        Inches(8.5), Inches(0.5),
        "  Onlayn imtihonlar uchun zamonaviy raqamli yechim",
        size=18, color=WHITE,
    )

    # Author block
    add_text(
        slide, Inches(0.7), Inches(6.3),
        Inches(8), Inches(0.4),
        "Muallif:  Bekzod Gulomov",
        size=14, color=WHITE,
    )
    add_text(
        slide, Inches(0.7), Inches(6.7),
        Inches(8), Inches(0.4),
        "Sana:  15-may, 2026-yil",
        size=14, color=WHITE,
    )

    # NSUMT badge bottom right
    add_round_rect(
        slide, prs.slide_width - Inches(2.5), prs.slide_height - Inches(1.1),
        Inches(1.9), Inches(0.55),
        WHITE, corner=0.5,
    )
    add_text(
        slide, prs.slide_width - Inches(2.5), prs.slide_height - Inches(1.0),
        Inches(1.9), Inches(0.4),
        "nsumt.uz",
        size=18, bold=True, color=BRAND, align=PP_ALIGN.CENTER,
    )
    return slide


def section_slide(prs, page_num, total, number, kicker, title):
    slide = blank_slide(prs)
    add_image_full(slide, ASSETS / "section_bg.png", prs)

    # Big number
    add_text(
        slide, Inches(0.8), Inches(1.5),
        Inches(4), Inches(2),
        number,
        size=180, bold=True, color=WHITE,
    )
    # Vertical accent line
    add_rect(
        slide, Inches(4.5), Inches(2.0),
        Inches(0.06), Inches(3.0),
        ACCENT,
    )
    add_text(
        slide, Inches(4.9), Inches(2.2),
        Inches(8), Inches(0.5),
        kicker.upper(),
        size=16, bold=True, color=ACCENT,
    )
    add_text(
        slide, Inches(4.9), Inches(2.8),
        Inches(8), Inches(2),
        title,
        size=48, bold=True, color=WHITE,
    )
    add_footer(slide, prs, page_num, total)
    return slide


def content_slide(prs, page_num, total, title, subtitle, bullets, icon_name):
    slide = blank_slide(prs)
    add_image_full(slide, ASSETS / "content_bg.png", prs)

    # Page number tag top-right
    add_round_rect(
        slide, prs.slide_width - Inches(1.3), Inches(0.4),
        Inches(0.8), Inches(0.4), BRAND, corner=0.5,
    )
    add_text(
        slide, prs.slide_width - Inches(1.3), Inches(0.45),
        Inches(0.8), Inches(0.3),
        f"{page_num:02d}",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )

    # Title
    add_text(
        slide, Inches(0.7), Inches(0.5),
        Inches(11), Inches(0.7),
        title, size=34, bold=True, color=BRAND_DARK,
    )
    if subtitle:
        add_text(
            slide, Inches(0.7), Inches(1.2),
            Inches(11), Inches(0.5),
            subtitle, size=16, color=GREY, italic=True,
        )

    # Accent underline
    add_rect(
        slide, Inches(0.7), Inches(1.78),
        Inches(0.8), Inches(0.06), ACCENT,
    )

    # Two-column: bullets left, icon right
    add_bullets(
        slide,
        Inches(0.9), Inches(2.2),
        Inches(8.5), Inches(4.5),
        bullets, size=18, color=DARK, line_spacing=14,
    )

    # Icon card on right
    icon_path = ASSETS / f"icon_{icon_name}.png"
    if icon_path.exists():
        slide.shapes.add_picture(
            str(icon_path),
            Inches(9.8), Inches(2.4),
            width=Inches(3.0), height=Inches(3.0),
        )
        # Decorative blob behind icon
        add_round_rect(
            slide, Inches(9.6), Inches(5.6),
            Inches(3.4), Inches(0.15),
            BRAND_LIGHT, corner=0.5,
        )

    add_footer(slide, prs, page_num, total)
    return slide


def closing_slide(prs, page_num, total):
    slide = blank_slide(prs)
    add_image_full(slide, ASSETS / "hero_title.png", prs)

    # Logo card centered top
    card_w = Inches(3)
    card_h = Inches(1.7)
    card_left = (prs.slide_width - card_w) // 2
    add_round_rect(
        slide, card_left, Inches(0.7), card_w, card_h,
        WHITE, corner=0.15,
    )
    slide.shapes.add_picture(
        str(LOGO),
        card_left + Inches(0.2), Inches(0.85),
        width=Inches(2.6),
    )

    add_text(
        slide, Inches(0), Inches(3.0),
        prs.slide_width, Inches(1.5),
        "Rahmat!",
        size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, Inches(0), Inches(4.4),
        prs.slide_width, Inches(0.6),
        "E'tiboringiz uchun tashakkur",
        size=24, color=LIGHT_BG, align=PP_ALIGN.CENTER,
    )

    # Accent bar centered
    bar_w = Inches(2.5)
    add_rect(
        slide, (prs.slide_width - bar_w) // 2, Inches(5.2),
        bar_w, Inches(0.08), ACCENT,
    )

    # Contact card
    contact_w = Inches(5)
    contact_h = Inches(1.3)
    contact_left = (prs.slide_width - contact_w) // 2
    add_round_rect(
        slide, contact_left, Inches(5.6),
        contact_w, contact_h, BRAND_DARK, corner=0.15,
    )
    add_text(
        slide, contact_left, Inches(5.75),
        contact_w, Inches(0.4),
        "Savollar va takliflar uchun",
        size=14, color=LIGHT_BG, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, contact_left, Inches(6.15),
        contact_w, Inches(0.5),
        "bgulomov000@gmail.com",
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, contact_left, Inches(6.7),
        contact_w, Inches(0.4),
        "nsumt.uz",
        size=14, color=LIGHT_BG, align=PP_ALIGN.CENTER,
    )
    return slide


SLIDES = [
    {
        "icon": "problem",
        "title": "Muammo",
        "subtitle": "Hozirgi vaziyat va qiyinchiliklar",
        "bullets": [
            "An'anaviy onlayn testlarda nazorat tizimi yo'q — ko'chirib yozish keng tarqalgan",
            "Talaba va o'qituvchi ma'lumotlari turli tizimlarda qo'lda kiritiladi",
            "Imtihon natijalari bo'yicha chuqur tahlil va statistika yetishmaydi",
            "Talabalarning psixologik holatini baholash uchun yagona platforma yo'q",
            "HEMIS bilan integratsiya cheklangan — ma'lumotlar ikki marta kiritiladi",
        ],
    },
    {
        "icon": "solution",
        "title": "Yechim",
        "subtitle": "Bizning yondashuvimiz",
        "bullets": [
            "Yagona platforma: testlar, psixologik baholash, talabalar boshqaruvi",
            "Sun'iy intellekt asosida real vaqtda yuz aniqlash orqali nazorat",
            "HEMIS tizimi bilan to'liq integratsiya — ma'lumotlar avtomatik sinxronizatsiya",
            "8 ta yo'nalishda batafsil statistika va analitika paneli",
            "Bulutli arxitektura: Docker konteynerlarida xavfsiz va kengaytiriladigan",
        ],
    },
    {
        "icon": "platform",
        "title": "Platforma haqida",
        "subtitle": "Umumiy ko'rinish",
        "bullets": [
            "NDKTU uchun maxsus ishlab chiqilgan ta'lim platformasi",
            "5 ta foydalanuvchi roli: administrator, o'qituvchi, talaba, psixolog, tyutor",
            "12 dan ortiq bog'langan modul: testlar, darslar, resurslar, baholash",
            "Veb-brauzerda ishlaydi — qo'shimcha dastur o'rnatish shart emas",
            "O'zbek va ingliz tillarida foydalanuvchi interfeysi",
        ],
    },
    {
        "icon": "users",
        "title": "Foydalanuvchilar va rollar",
        "subtitle": "Har bir rol — o'z imkoniyatlari",
        "bullets": [
            "Administrator: tizimni boshqarish, foydalanuvchilar, rollar va ruxsatlar",
            "O'qituvchi: testlar yaratish, darslarni boshqarish, natijalarni ko'rish",
            "Talaba: testlarda qatnashish, psixologik baholash, natijalarni ko'rish",
            "Psixolog: psixologik metodikalar yaratish, talabalar holatini kuzatish",
            "Tyutor: guruhlarni nazorat qilish, talabalarga akademik yordam",
        ],
    },
    {
        "icon": "modules",
        "title": "Asosiy modullar",
        "subtitle": "Platforma tarkibiy qismlari",
        "bullets": [
            "Test tizimi (Quiz) — PIN kod, taymer, avtomatik baholash",
            "Yuz aniqlash (Face Detection) — real vaqtda imtihon nazorati",
            "Psixologik baholash — standart metodikalar va avtomatik skoring",
            "HEMIS integratsiyasi — talaba va o'qituvchilar avtomatik import",
            "Statistika va tahlil — 8 ta bo'lim bo'yicha analitik panel",
            "Tashkiliy boshqaruv — fakultet, kafedra, guruh, o'quv yili",
        ],
    },
    {
        "icon": "quiz",
        "title": "Test tizimi",
        "subtitle": "Onlayn imtihonlarni tashkil etish",
        "bullets": [
            "O'qituvchilar test yaratadi: bir nechta variantli, ha/yo'q, shkala, ochiq savol",
            "Talabalar PIN kod orqali testga kiradi — ruxsatsiz kirish to'sib qo'yilgan",
            "Real vaqt taymeri — vaqt tugashi bilan test avtomatik yopiladi",
            "Obyektiv savollar avtomatik baholanadi, natija darhol ko'rinadi",
            "Testni nusxalash, qayta o'tkazish va savol bankidan foydalanish",
        ],
    },
    {
        "icon": "face",
        "title": "Yuz aniqlash (Face Detection)",
        "subtitle": "Sun'iy intellekt asosida nazorat",
        "bullets": [
            "WebSocket orqali brauzerdan jonli video oqim olinadi",
            "Google MediaPipe modeli har bir kadrda yuzlar sonini aniqlaydi",
            "Ikkita yuz aniqlanganda yoki shaxs o'zgarganda ogohlantirish beriladi",
            "3 ta ogohlantirishdan keyin test avtomatik tugatiladi",
            "Buzilishlar dalil sifatida rasm shaklida saqlanadi",
        ],
    },
    {
        "icon": "psychology",
        "title": "Psixologik baholash",
        "subtitle": "Talabalar ruhiy holatini kuzatish",
        "bullets": [
            "Psixologlar standart metodikalar yaratadi (Likert, ha/yo'q, ko'p variantli)",
            "Talabalar onlayn rejimda psixologik testlarni topshiradi",
            "Tizim natijalarni avtomatik hisoblaydi va talqin qiladi",
            "Rasm orqali stimul beruvchi savollar qo'llab-quvvatlanadi",
            "Talabalar bo'yicha psixologik tarix saqlanadi va kuzatib boriladi",
        ],
    },
    {
        "icon": "org",
        "title": "Tashkiliy boshqaruv",
        "subtitle": "Universitet strukturasini raqamlashtirish",
        "bullets": [
            "Fakultet → Kafedra → Guruh → Talaba — yagona ierarxiya",
            "O'quv yili va semestrlar bo'yicha rejalashtirish",
            "Sinflar (mashg'ulotlar) o'qituvchi, guruh va mavzu bilan bog'lanadi",
            "Davomat va talabalarning aktivligini kuzatish",
            "Excel orqali ommaviy import va eksport qilish imkoniyati",
        ],
    },
    {
        "icon": "content",
        "title": "O'quv kontenti",
        "subtitle": "Darslar va resurslar",
        "bullets": [
            "Darslar: mavzu, tavsif, biriktirilgan o'qituvchi va guruh",
            "Resurslar: PDF, hujjat va rasm shaklidagi o'quv materiallari",
            "Syllabus: kurs dasturini elektron shaklda yuritish",
            "Yakuniy imtihonlar: oraliq testlardan alohida hisobga olinadi",
            "Mavzular bo'yicha boy matn muharriri (rasm va jadvallar bilan)",
        ],
    },
    {
        "icon": "hemis",
        "title": "HEMIS integratsiyasi",
        "subtitle": "Rasmiy ta'lim tizimi bilan bog'lanish",
        "bullets": [
            "Administrator HEMIS login orqali bevosita ulanadi",
            "Talabalar, o'qituvchilar, fanlar va guruhlar avtomatik import qilinadi",
            "Sinxronlashdan oldin oldindan ko'rib chiqish (preview) imkoniyati",
            "Har bir sinxronlash tranzaksiyasi tarixda saqlanadi",
            "Qo'lda kiritish zaruriyati yo'qoladi — vaqt va xato kamayadi",
        ],
    },
    {
        "icon": "stats",
        "title": "Statistika va tahlil",
        "subtitle": "Ma'lumotlarga asoslangan qarorlar",
        "bullets": [
            "Umumiy statistika — tizim bo'yicha asosiy ko'rsatkichlar",
            "Test analitikasi — har bir test bo'yicha o'tish foizi va vaqt",
            "Savol tahlili (Item Analysis) — savol qiyinligi va sifati",
            "Demografiya — fakultet va guruhlar bo'yicha tahlil",
            "Nazorat statistikasi — buzilishlar sababi va miqyosi",
            "Psixologik, o'qituvchi faolligi va yakuniy imtihon hisobotlari",
        ],
    },
    {
        "icon": "tech",
        "title": "Texnologiyalar",
        "subtitle": "Zamonaviy va ishonchli stek",
        "bullets": [
            "Backend: FastAPI (Python 3.12) — yuqori unumdor async API",
            "Frontend: React 19 + TypeScript + Tailwind CSS",
            "Ma'lumotlar bazasi: PostgreSQL 17 + Redis (keshlash)",
            "Yuz aniqlash: Google MediaPipe + OpenCV",
            "Autentifikatsiya: JWT tokenlari + bcrypt parol xeshlash",
            "Konteynerlash: Docker + Docker Compose",
        ],
    },
    {
        "icon": "arch",
        "title": "Arxitektura",
        "subtitle": "Mikroservislar yondashuvi",
        "bullets": [
            "Frontend (port 3000) — React SPA, Nginx orqali xizmat ko'rsatadi",
            "Backend API (port 8000) — asosiy biznes mantiq va REST API",
            "Face Detection xizmati (port 8001) — alohida mikroservis",
            "PostgreSQL + Redis — ma'lumotlar va kesh qatlami",
            "Health check va avtomatik qayta ishga tushirish mexanizmi",
            "Production-da Nginx reverse proxy orqali xavfsiz ulanish",
        ],
    },
    {
        "icon": "security",
        "title": "Xavfsizlik",
        "subtitle": "Ma'lumotlar himoyasi va yaxlitligi",
        "bullets": [
            "JWT tokenlar va xavfsiz autentifikatsiya — sessiyalar himoyalangan",
            "Rollar va ruxsatlar (RBAC) — har bir amal nazorat ostida",
            "Rate limiting — brute-force va DDoS hujumlariga qarshi himoya",
            "XSS himoyasi — barcha HTML kontent DOMPurify orqali tozalanadi",
            "SQL injection himoyasi — SQLAlchemy parametrlangan so'rovlari",
            "Imtihon buzilishlari dalillari saqlanadi — to'liq audit izi",
        ],
    },
    {
        "icon": "future",
        "title": "Kelajakdagi imkoniyatlar",
        "subtitle": "Rivojlanish yo'nalishlari",
        "bullets": [
            "Mobil ilova (iOS / Android) — talabalar uchun qulaylik",
            "Sun'iy intellekt yordamida savol generatsiyasi (AI)",
            "Moslashuvchan testlash (CAT) — daraja bo'yicha avtomatik moslashish",
            "Kengaytirilgan nazorat: ko'z harakati va klaviatura tahlili",
            "LMS tizimlari bilan integratsiya (Moodle, Canvas, Blackboard)",
            "Prognozli analitika — qiyinchilik bo'lishi mumkin talabalarni aniqlash",
        ],
    },
    {
        "icon": "conclusion",
        "title": "Xulosa",
        "subtitle": "Platforma universitetga nima beradi",
        "bullets": [
            "Imtihon jarayonini to'liq raqamlashtirish va shaffoflashtirish",
            "Ko'chirib yozish va firibgarliklarning keskin kamayishi (AI nazorat)",
            "Ma'lumotlarni qo'lda kiritish vaqtining qisqarishi (HEMIS sinxronlash)",
            "O'qituvchilar va rahbariyat uchun chuqur analitika",
            "Talabalar psixologik holatini kuzatish va o'z vaqtida yordam berish",
            "Kelajakda kengaytirish uchun moslashuvchan zamonaviy arxitektura",
        ],
    },
]


def main():
    if not LOGO.exists():
        raise FileNotFoundError(f"Logo not found: {LOGO}")
    if not (ASSETS / "hero_title.png").exists():
        raise FileNotFoundError(
            f"Assets missing. Run: uv run --with pillow python scripts/generate_assets.py"
        )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(SLIDES) + 2

    title_slide(prs, total)

    for idx, data in enumerate(SLIDES, start=2):
        content_slide(
            prs, idx, total,
            data["title"], data["subtitle"], data["bullets"],
            data["icon"],
        )

    closing_slide(prs, total, total)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {total}")


if __name__ == "__main__":
    main()
